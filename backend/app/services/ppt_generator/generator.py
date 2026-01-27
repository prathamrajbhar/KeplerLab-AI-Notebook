import os
from io import BytesIO
from typing import Dict, Optional, Callable, List

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .html_renderer import HTMLSlideRenderer


class Colors:
    BG_PRIMARY = RGBColor(15, 17, 23)
    BG_SECONDARY = RGBColor(22, 27, 34)
    BG_TERTIARY = RGBColor(30, 36, 45)
    BG_GLASS = RGBColor(35, 42, 52)
    ACCENT_PRIMARY = RGBColor(99, 102, 241)
    ACCENT_SECONDARY = RGBColor(139, 92, 246)
    ACCENT_CYAN = RGBColor(34, 211, 238)
    TEXT_PRIMARY = RGBColor(248, 250, 252)
    TEXT_SECONDARY = RGBColor(203, 213, 225)
    TEXT_MUTED = RGBColor(148, 163, 184)
    GLASS_BORDER = RGBColor(71, 85, 105)


class LayoutConfig:
    """Configuration for slide layout constraints"""
    CONTENT_TOP = 1.2
    CONTENT_LEFT = 0.5
    CONTENT_WIDTH = 12.3  # Full width minus margins
    CONTENT_BOTTOM = 6.6
    BLOCK_GAP = 0.25  # More space between blocks
    
    # Font sizes (FIXED - no dynamic scaling)
    TITLE_FONT_SIZE = 28
    HEADING_FONT_SIZE = 16
    PARA_FONT_SIZE = 14
    BULLET_FONT_SIZE = 13
    TABLE_HEADER_SIZE = 11
    TABLE_CELL_SIZE = 10
    DEFINITION_TERM_SIZE = 16
    DEFINITION_TEXT_SIZE = 13
    
    # Character limits
    MAX_PARA_CHARS = 450
    MAX_BULLET_CHARS = 150
    MAX_BULLETS = 7
    MAX_POINTS = 7

    @classmethod
    def available_height(cls):
        return cls.CONTENT_BOTTOM - cls.CONTENT_TOP


def truncate_text(text: str, max_chars: int) -> str:
    """Truncate text to max characters with ellipsis"""
    if not text or len(text) <= max_chars:
        return text
    return text[:max_chars-3].rsplit(' ', 1)[0] + '...'


def estimate_text_height(text: str, width_inches: float, font_size: int) -> float:
    """Estimate height needed for text in inches"""
    if not text:
        return 0
    chars_per_line = int(width_inches * 12)
    lines = max(1, len(text) // chars_per_line + 1)
    line_height = font_size / 72 * 1.4
    return lines * line_height


def clean_text(text: str) -> str:
    """Remove markdown formatting"""
    if not text:
        return ''
    return text.replace('**', '').replace('*', '').replace('`', '').strip()


class PPTGenerator:
    def __init__(self, image_generator: Optional[Callable] = None):
        self.image_generator = image_generator
        self.html_renderer = HTMLSlideRenderer(image_generator)
    
    def generate(self, chapter_data: Dict) -> BytesIO:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        chapter_title = chapter_data.get('chapter_title', 'Presentation')
        slides = chapter_data.get('slides', [])
        
        self._add_title_slide(prs, chapter_title)
        
        for i, slide_data in enumerate(slides):
            self._add_content_slide(prs, slide_data, i + 2)
        
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer
    
    def save_html(self, chapter_data: Dict, output_dir: str = "output/html") -> str:
        os.makedirs(output_dir, exist_ok=True)
        html = self.html_renderer.render_presentation(chapter_data)
        title = chapter_data.get('chapter_title', 'presentation')
        safe_title = "".join(c if c.isalnum() or c in (' ', '-', '_') else '' for c in title)
        safe_title = safe_title.replace(' ', '_')[:50] or 'presentation'
        filepath = os.path.join(output_dir, f"{safe_title}.html")
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(html)
        return filepath
    
    def _add_background(self, slide, color=None):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color or Colors.BG_PRIMARY
    
    def _add_title_slide(self, prs: Presentation, title: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_background(slide)
        
        # Right panel
        rp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5), Inches(0), Inches(3.833), Inches(7.5))
        rp.fill.solid()
        rp.fill.fore_color.rgb = Colors.BG_SECONDARY
        rp.line.fill.background()
        
        # Bottom bar
        bb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15))
        bb.fill.solid()
        bb.fill.fore_color.rgb = Colors.ACCENT_PRIMARY
        bb.line.fill.background()
        
        # Left accent
        la = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.2), Inches(0.08), Inches(2.5))
        la.fill.solid()
        la.fill.fore_color.rgb = Colors.ACCENT_CYAN
        la.line.fill.background()
        
        # Decorative circles
        for x, y, s, c in [(10.5, 0.6, 0.5, Colors.ACCENT_PRIMARY), 
                           (11.2, 0.9, 0.35, Colors.ACCENT_SECONDARY),
                           (11.7, 0.5, 0.25, Colors.ACCENT_CYAN)]:
            circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(s), Inches(s))
            circ.fill.solid()
            circ.fill.fore_color.rgb = c
            circ.line.fill.background()
        
        # Title
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8), Inches(2.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = Colors.TEXT_PRIMARY
        p.font.name = "Segoe UI"
        
        # Subtitle
        sb = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(4), Inches(0.5))
        p = sb.text_frame.paragraphs[0]
        p.text = "PRESENTATION OVERVIEW"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = Colors.ACCENT_CYAN
        p.font.name = "Segoe UI"
        
        # Underline
        ul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.6), Inches(2.5), Inches(0.04))
        ul.fill.solid()
        ul.fill.fore_color.rgb = Colors.ACCENT_CYAN
        ul.line.fill.background()
    
    def _add_content_slide(self, prs: Presentation, slide_data: Dict, page_num: int):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_background(slide)
        
        # Top accent bar (full width)
        tb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.05))
        tb.fill.solid()
        tb.fill.fore_color.rgb = Colors.ACCENT_PRIMARY
        tb.line.fill.background()
        
        section = slide_data.get('section_number', '')
        title = slide_data.get('title', '')
        blocks = slide_data.get('blocks', [])
        
        full_title = f"{section}   {title}" if section else title
        
        # Title (full width)
        ttb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.75))
        tf = ttb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = truncate_text(full_title, 100)
        p.font.size = Pt(LayoutConfig.TITLE_FONT_SIZE)
        p.font.bold = True
        p.font.color.rgb = Colors.TEXT_PRIMARY
        p.font.name = "Segoe UI"
        
        # Title underline (wider)
        tul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(5), Inches(0.04))
        tul.fill.solid()
        tul.fill.fore_color.rgb = Colors.ACCENT_CYAN
        tul.line.fill.background()
        
        # Accent dot
        dt = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.6), Inches(0.93), Inches(0.08), Inches(0.08))
        dt.fill.solid()
        dt.fill.fore_color.rgb = Colors.ACCENT_PRIMARY
        dt.line.fill.background()
        
        # Footer
        fb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4))
        fb.fill.solid()
        fb.fill.fore_color.rgb = Colors.BG_TERTIARY
        fb.line.fill.background()
        
        # Page number
        pb = slide.shapes.add_textbox(Inches(12.3), Inches(7.15), Inches(0.8), Inches(0.3))
        p = pb.text_frame.paragraphs[0]
        p.text = f"{page_num:02d}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = Colors.TEXT_MUTED
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.RIGHT
        
        # Render blocks with proper spacing
        current_top = Inches(LayoutConfig.CONTENT_TOP)
        left = Inches(LayoutConfig.CONTENT_LEFT)
        width = Inches(LayoutConfig.CONTENT_WIDTH)
        max_bottom = Inches(LayoutConfig.CONTENT_BOTTOM)
        
        for block in blocks:
            if current_top >= max_bottom:
                break
            current_top = self._render_block(slide, block, left, current_top, width, max_bottom)
    
    def _render_block(self, slide, block: Dict, left, top, width, max_bottom):
        block_type = block.get('block_type', '')
        
        renderers = {
            'paragraph_block': self._render_paragraph,
            'bullet_block': self._render_bullets,
            'numbered_block': self._render_numbered,
            'definition_block': self._render_definition,
            'table_block': self._render_table,
            'image_block': self._render_image_block,
            'diagram_block': self._render_diagram_block,
        }
        
        renderer = renderers.get(block_type)
        if renderer:
            return renderer(slide, block, left, top, width, max_bottom)
        return top
    
    def _render_paragraph(self, slide, block: Dict, left, top, width, max_bottom):
        content = block.get('content', '')
        if not content:
            return top
        
        content = truncate_text(clean_text(content), LayoutConfig.MAX_PARA_CHARS)
        
        # Calculate lines needed
        chars_per_line = int((width / Inches(1)) * 10)
        lines = max(1, len(content) // chars_per_line + 1)
        height = Inches(min(lines * 0.22 + 0.1, 1.8))
        
        if top + height > max_bottom:
            height = max_bottom - top
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(LayoutConfig.PARA_FONT_SIZE)
        p.font.color.rgb = Colors.TEXT_SECONDARY
        p.font.name = "Segoe UI"
        p.line_spacing = 1.4
        
        return top + height + Inches(LayoutConfig.BLOCK_GAP)
    
    def _render_bullets(self, slide, block: Dict, left, top, width, max_bottom):
        heading = block.get('heading', '')
        bullets = block.get('bullets', [])
        if not bullets:
            return top
        
        item_height = 0.35
        
        if heading:
            if top + Inches(0.4) > max_bottom:
                return top
            hBox = slide.shapes.add_textbox(left, top, width, Inches(0.35))
            tf = hBox.text_frame
            p = tf.paragraphs[0]
            p.text = truncate_text(heading, 80)
            p.font.size = Pt(LayoutConfig.HEADING_FONT_SIZE)
            p.font.bold = True
            p.font.color.rgb = Colors.TEXT_PRIMARY
            p.font.name = "Segoe UI"
            top += Inches(0.38)
        
        for bullet in bullets[:LayoutConfig.MAX_BULLETS]:
            if top + Inches(item_height) > max_bottom:
                break
            
            # Diamond bullet
            diamond = slide.shapes.add_shape(
                MSO_SHAPE.DIAMOND, 
                left + Inches(0.05), top + Inches(0.08), 
                Inches(0.12), Inches(0.12)
            )
            diamond.fill.solid()
            diamond.fill.fore_color.rgb = Colors.ACCENT_CYAN
            diamond.line.fill.background()
            
            # Text
            bullet_text = truncate_text(clean_text(bullet), LayoutConfig.MAX_BULLET_CHARS)
            txBox = slide.shapes.add_textbox(left + Inches(0.28), top, width - Inches(0.3), Inches(item_height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = bullet_text
            p.font.size = Pt(LayoutConfig.BULLET_FONT_SIZE)
            p.font.color.rgb = Colors.TEXT_SECONDARY
            p.font.name = "Segoe UI"
            
            top += Inches(item_height)
        
        return top + Inches(LayoutConfig.BLOCK_GAP)
    
    def _render_numbered(self, slide, block: Dict, left, top, width, max_bottom):
        heading = block.get('heading', '')
        points = block.get('points', [])
        if not points:
            return top
        
        item_height = 0.38
        
        if heading:
            if top + Inches(0.4) > max_bottom:
                return top
            hBox = slide.shapes.add_textbox(left, top, width, Inches(0.35))
            tf = hBox.text_frame
            p = tf.paragraphs[0]
            p.text = truncate_text(heading, 80)
            p.font.size = Pt(LayoutConfig.HEADING_FONT_SIZE)
            p.font.bold = True
            p.font.color.rgb = Colors.TEXT_PRIMARY
            p.font.name = "Segoe UI"
            top += Inches(0.38)
        
        for i, point in enumerate(points[:LayoutConfig.MAX_POINTS]):
            if top + Inches(item_height) > max_bottom:
                break
            
            # Number circle
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Inches(0.04), Inches(0.26), Inches(0.26))
            circle.fill.solid()
            circle.fill.fore_color.rgb = Colors.ACCENT_PRIMARY
            circle.line.fill.background()
            
            # Number
            numBox = slide.shapes.add_textbox(left, top + Inches(0.05), Inches(0.26), Inches(0.24))
            tf = numBox.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = Colors.TEXT_PRIMARY
            p.font.name = "Segoe UI"
            p.alignment = PP_ALIGN.CENTER
            
            # Text
            point_text = truncate_text(clean_text(point), LayoutConfig.MAX_BULLET_CHARS)
            txBox = slide.shapes.add_textbox(left + Inches(0.36), top, width - Inches(0.4), Inches(item_height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = point_text
            p.font.size = Pt(LayoutConfig.BULLET_FONT_SIZE)
            p.font.color.rgb = Colors.TEXT_SECONDARY
            p.font.name = "Segoe UI"
            
            top += Inches(item_height)
        
        return top + Inches(LayoutConfig.BLOCK_GAP)
    
    def _render_definition(self, slide, block: Dict, left, top, width, max_bottom):
        term = block.get('term', '')
        definition = block.get('definition', '')
        if not term and not definition:
            return top
        
        # Calculate box height based on content
        def_text = truncate_text(clean_text(definition), 350) if definition else ''
        lines = len(def_text) // 120 + 1 if def_text else 0
        box_height = 0.5 + (lines * 0.22) if def_text else 0.45
        box_height = min(box_height, 1.4)
        
        if top + Inches(box_height) > max_bottom:
            return top
        
        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(box_height))
        bg.fill.solid()
        bg.fill.fore_color.rgb = Colors.BG_GLASS
        bg.line.color.rgb = Colors.GLASS_BORDER
        bg.line.width = Pt(1)
        
        # Accent bar
        ab = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.12), top + Inches(0.12), Inches(0.05), Inches(box_height - 0.24))
        ab.fill.solid()
        ab.fill.fore_color.rgb = Colors.ACCENT_CYAN
        ab.line.fill.background()
        
        if term:
            tBox = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.1), width - Inches(0.5), Inches(0.32))
            p = tBox.text_frame.paragraphs[0]
            p.text = truncate_text(term, 70)
            p.font.size = Pt(LayoutConfig.DEFINITION_TERM_SIZE)
            p.font.bold = True
            p.font.color.rgb = Colors.TEXT_PRIMARY
            p.font.name = "Segoe UI"
        
        if def_text:
            dBox = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.42), width - Inches(0.5), Inches(box_height - 0.52))
            tf = dBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = def_text
            p.font.size = Pt(LayoutConfig.DEFINITION_TEXT_SIZE)
            p.font.color.rgb = Colors.TEXT_SECONDARY
            p.font.name = "Segoe UI"
            p.line_spacing = 1.35
        
        return top + Inches(box_height) + Inches(LayoutConfig.BLOCK_GAP)
    
    def _render_table(self, slide, block: Dict, left, top, width, max_bottom):
        heading = block.get('heading', '')
        columns = block.get('columns', [])
        rows = block.get('rows', [])
        if not columns or not rows:
            return top
        
        if heading:
            if top + Inches(0.35) > max_bottom:
                return top
            hBox = slide.shapes.add_textbox(left, top, width, Inches(0.32))
            p = hBox.text_frame.paragraphs[0]
            p.text = truncate_text(heading, 80)
            p.font.size = Pt(LayoutConfig.HEADING_FONT_SIZE)
            p.font.bold = True
            p.font.color.rgb = Colors.TEXT_PRIMARY
            p.font.name = "Segoe UI"
            top += Inches(0.35)
        
        num_cols = min(len(columns), 6)
        row_height = 0.35
        num_rows = min(len(rows) + 1, 7)
        table_height = row_height * num_rows
        
        if top + Inches(table_height) > max_bottom:
            available_rows = int((max_bottom - top) / Inches(row_height)) - 1
            if available_rows < 2:
                return top
            num_rows = available_rows + 1
            table_height = row_height * num_rows
        
        table = slide.shapes.add_table(num_rows, num_cols, left, top, width, Inches(table_height)).table
        
        # Header
        for i, col in enumerate(columns[:num_cols]):
            cell = table.cell(0, i)
            cell.text = str(col)[:30]
            cell.fill.solid()
            cell.fill.fore_color.rgb = Colors.ACCENT_PRIMARY
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(LayoutConfig.TABLE_HEADER_SIZE)
            p.font.bold = True
            p.font.color.rgb = Colors.TEXT_PRIMARY
            p.font.name = "Segoe UI"
            p.alignment = PP_ALIGN.CENTER
        
        # Data rows
        for r, row in enumerate(rows[:num_rows-1]):
            for c, val in enumerate(row[:num_cols]):
                cell = table.cell(r + 1, c)
                cell.text = str(val)[:45]
                cell.fill.solid()
                cell.fill.fore_color.rgb = Colors.BG_GLASS if r % 2 == 0 else Colors.BG_SECONDARY
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(LayoutConfig.TABLE_CELL_SIZE)
                p.font.color.rgb = Colors.TEXT_SECONDARY
                p.font.name = "Segoe UI"
        
        return top + Inches(table_height + 0.2)
    
    def _render_image_block(self, slide, block: Dict, left, top, width, max_bottom):
        prompt = block.get('image_prompt', '')
        if not prompt:
            return top
            
        image_data = None
        if self.image_generator:
            try:
                image_data = self.image_generator(prompt)
            except Exception:
                pass
                
        if image_data:
            img_height = min(3.0, (max_bottom - top) / Inches(1) - 0.2)
            if img_height > 0.5:
                slide.shapes.add_picture(BytesIO(image_data), left, top, width=width, height=Inches(img_height))
                return top + Inches(img_height + LayoutConfig.BLOCK_GAP)
        
        # Placeholder
        txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        p = txBox.text_frame.paragraphs[0]
        p.text = f"[Image: {prompt[:60]}...]"
        p.font.size = Pt(11)
        p.font.color.rgb = Colors.TEXT_MUTED
        return top + Inches(0.5)

    def _render_diagram_block(self, slide, block: Dict, left, top, width, max_bottom):
        prompt = block.get('diagram_prompt', '')
        if not prompt:
            return top
            
        image_data = None
        if self.image_generator:
            try:
                query = prompt if "diagram" in prompt.lower() else prompt + " diagram"
                image_data = self.image_generator(query)
            except Exception:
                pass
                
        if image_data:
            img_height = min(3.0, (max_bottom - top) / Inches(1) - 0.2)
            if img_height > 0.5:
                slide.shapes.add_picture(BytesIO(image_data), left, top, width=width, height=Inches(img_height))
                return top + Inches(img_height + LayoutConfig.BLOCK_GAP)
        
        # Placeholder
        txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        p = txBox.text_frame.paragraphs[0]
        p.text = f"[Diagram: {prompt[:60]}...]"
        p.font.size = Pt(11)
        p.font.color.rgb = Colors.TEXT_MUTED
        return top + Inches(0.5)


def generate_ppt(chapter_data: Dict, image_generator=None) -> BytesIO:
    return PPTGenerator(image_generator).generate(chapter_data)
